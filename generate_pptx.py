import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Color Palette: Slate & Sky & Amber (Modern Dark Mode)
    BG_COLOR = RGBColor(15, 23, 42)        # Slate 900
    CARD_BG = RGBColor(30, 41, 59)        # Slate 800
    BORDER_COLOR = RGBColor(71, 85, 105)   # Slate 600
    
    TEXT_MAIN = RGBColor(255, 255, 255)    # White
    TEXT_MUTED = RGBColor(203, 213, 225)   # Slate 300
    
    ACCENT_SKY = RGBColor(56, 189, 248)    # Sky 400
    ACCENT_AMBER = RGBColor(251, 191, 36)  # Amber 400

    blank_layout = prs.slide_layouts[6]

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title, category=""):
        # Category/Tag Tracker
        if category:
            cat_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.3))
            cat_tf = cat_box.text_frame
            cat_tf.word_wrap = True
            cat_p = cat_tf.paragraphs[0]
            cat_p.text = category.upper()
            cat_p.font.name = "Segoe UI"
            cat_p.font.size = Pt(11)
            cat_p.font.bold = True
            cat_p.font.color.rgb = ACCENT_SKY
        
        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.7))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.name = "Segoe UI"
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = TEXT_MAIN

    def add_screenshot_placeholder(slide, left, top, width, height, label):
        # Outer Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        
        # Inner dashed dashed border for photo
        inner_w = width - Inches(0.4)
        inner_h = height - Inches(0.4)
        inner_l = left + Inches(0.2)
        inner_t = top + Inches(0.2)
        
        inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inner_l, inner_t, inner_w, inner_h)
        inner.fill.solid()
        inner.fill.fore_color.rgb = BG_COLOR
        inner.line.color.rgb = ACCENT_SKY
        inner.line.width = Pt(1)
        
        # Text label inside inner box
        tf = inner.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[ PLACE PHOTO HERE ]\n\n{label}"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.font.bold = True

    def add_bullets(slide, left, top, width, height, bullets):
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = b
            p.font.name = "Segoe UI"
            p.font.size = Pt(15)
            p.font.color.rgb = TEXT_MUTED
            p.level = 0
            p.space_after = Pt(12)
            
            # Make first few words bold for scanability
            if ":" in b:
                parts = b.split(":", 1)
                p.text = parts[0] + ":"
                p.font.bold = True
                p.font.color.rgb = ACCENT_AMBER
                
                run = p.add_run()
                run.text = parts[1]
                run.font.bold = False
                run.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 1: Opportunity & Differentiation
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_background(s1)
    add_header(s1, "Vetting Exoplanets with Explainable AI", "Exoplanet Vetting Platform")
    
    s1_bullets = [
        "The Core Problem: Exoplanet vetting is a major research bottleneck. Deep learning classifiers fail silently due to data domain shifts.",
        "Physics-Guided Hybrid AI: We combine the AstroNet CNN model directly with 10 diagnostic physical features inside a multi-class Random Forest.",
        "Explainable AI (USP): Our SHAP explainability layer translates raw predictions into clear physical reasons (e.g. secondary eclipses) for astronomers."
    ]
    add_bullets(s1, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), s1_bullets)
    add_screenshot_placeholder(s1, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.8), "Dashboard Core Metrics View\nShows Period, Duration, SNR, and Depth fields.")

    # -------------------------------------------------------------
    # SLIDE 2: Core Platform Features
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_background(s2)
    add_header(s2, "Automated Detection & Processing", "Core Features")
    
    s2_bullets = [
        "MAST Archive Integration: Input any Kepler or TESS ID to fetch light curves automatically.",
        "3-Stage Detrending Pipeline: Removes outlier flares, stellar variability, and high-frequency noise.",
        "Box Least Squares (BLS): Scans 20,000 trial periods to identify periodic transit dips.",
        "Astrophysical Diagnostics: Auto-calculates transit shape (U-vs-V), odd-even depth variation, and symmetry."
    ]
    add_bullets(s2, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), s2_bullets)
    add_screenshot_placeholder(s2, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.8), "Light Curve Visualizer View\nShows detrended baseline and phase-folded transit dips.")

    # -------------------------------------------------------------
    # SLIDE 3: Process Flow Diagram
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_background(s3)
    add_header(s3, "Pipeline Process Flow", "Execution Pipeline")
    
    flow_steps = [
        ("1. Input KIC/TIC ID", "User searches star ID"),
        ("2. Data Download", "Fetch from NASA MAST"),
        ("3. Signal Cleansing", "3-Step Detrending"),
        ("4. Transit Search", "Box Least Squares"),
        ("5. Feature Extract", "10 Physical Metrics"),
        ("6. AI Classification", "AstroNet CNN + RF"),
        ("7. SHAP XAI Engine", "Explain contributions"),
        ("8. Interactive UI", "Visual results & PDF")
    ]
    
    # Render flow steps as horizontal cards with arrows
    for idx, (title, desc) in enumerate(flow_steps):
        row = idx // 4
        col = idx % 4
        
        left = Inches(0.5 + col * 3.1)
        top = Inches(2.0 + row * 2.3)
        width = Inches(2.8)
        height = Inches(1.6)
        
        # Step shape
        shape = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = ACCENT_SKY
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_AMBER
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 4: Interactive Vetting Dashboard
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_background(s4)
    add_header(s4, "Interactive Vetting Dashboard", "UI / Mockup Overview")
    
    desc_box = slide_desc = s4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.5))
    desc_p = desc_box.text_frame.paragraphs[0]
    desc_p.text = "A single unified portal for exoplanet researchers, integrating raw light curves, neural net predictions, and model explanations."
    desc_p.font.name = "Segoe UI"
    desc_p.font.size = Pt(15)
    desc_p.font.color.rgb = TEXT_MUTED
    
    add_screenshot_placeholder(s4, Inches(0.5), Inches(2.0), Inches(12.33), Inches(4.8), "Full-Width Frontend Dashboard Screenshot\nCapture the entire dashboard showing the curves, class indicators, and SHAP charts.")

    # -------------------------------------------------------------
    # SLIDE 5: System Architecture
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_background(s5)
    add_header(s5, "System Architecture & Integration", "Architecture")
    
    # We will draw 4 main component cards representing layers
    layers = [
        ("Presentation Layer", "React 19 Frontend\nTypeScript & Tailwind\nInteractive Recharts\nDynamic PDF Exporter", Inches(0.5), Inches(2.2), Inches(2.6), Inches(4.0)),
        ("Web API Gateway", "Tornado REST Server\nAsynchronous endpoints\nBase64 Plot Encoder\nJSON Marshalling", Inches(3.7), Inches(2.2), Inches(2.6), Inches(4.0)),
        ("Astrophysical Engine", "Lightkurve MAST Client\nSciPy Detrending Layer\nAstropy BLS Periodogram\nFeature Extractor Module", Inches(6.9), Inches(2.2), Inches(2.6), Inches(4.0)),
        ("ML & Explainability", "TensorFlow AstroNet CNN\nRule-Assisted RF Classifier\nSHAP TreeExplainer\nPDF Report Generator", Inches(10.1), Inches(2.2), Inches(2.6), Inches(4.0))
    ]
    
    for title, desc, l, t, w, h in layers:
        shape = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = BORDER_COLOR
        shape.line.width = Pt(2)
        
        tf = shape.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_AMBER
        p1.space_after = Pt(14)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 6: Technology Stack
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_background(s6)
    add_header(s6, "Scientific Technology Stack", "Tech Stack")
    
    s6_bullets_left = [
        "Data Source: Lightkurve API / Mikulski Space Telescope Archive (MAST).",
        "Physics & Analytics: Astropy (Box Least Squares), SciPy (Signal detrending).",
        "Web API Server: Tornado (High-concurrency Python REST server)."
    ]
    
    s6_bullets_right = [
        "Machine Learning: TensorFlow 1.x (AstroNet CNN), Scikit-Learn (Random Forest).",
        "Explainability & Docs: SHAP Explainer, ReportLab (PDF compiler).",
        "Frontend Interface: React 19, TypeScript, TailwindCSS, Recharts."
    ]
    
    add_bullets(s6, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), s6_bullets_left)
    add_bullets(s6, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0), s6_bullets_right)

    # -------------------------------------------------------------
    # SLIDE 7: Cost & Feasibility
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_background(s7)
    add_header(s7, "Estimated Cost & Scaling", "Cost & Feasibility")
    
    s7_bullets = [
        "Infrastructure Cost: Model inference runs on standard CPUs; no expensive GPU hosts required.",
        "Deployment: Frontend is hosted on Vercel (free tier). Backend container is hosted on Render ($7-$35/mo).",
        "Data Costs: $0.00. Free public API access to NASA MAST databases.",
        "Licensing: 100% open-source software stack with no licensing fees."
    ]
    add_bullets(s7, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), s7_bullets)
    
    # Add Cost Table
    rows = 5
    cols = 4
    table_shape = s7.shapes.add_table(rows, cols, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.5))
    table = table_shape.table
    
    # Headers
    headers = ["Category", "Tool / Service", "Dev Cost", "Scale Cost"]
    for c_idx, text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_SKY
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BG_COLOR
        
    data = [
        ["Data Access", "NASA MAST", "$0.00", "$0.00"],
        ["Hosting", "Vercel & Render", "$0.00", "$7 - $35/mo"],
        ["Libraries", "TensorFlow / Sklearn", "$0.00", "$0.00"],
        ["Total", "Platform Stack", "$0.00", "$7 - $35/mo"]
    ]
    
    for r_idx, row_vals in enumerate(data):
        for c_idx, val in enumerate(row_vals):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Segoe UI"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_MAIN
            if r_idx == 3: # Total Highlight
                p.font.bold = True
                p.font.color.rgb = ACCENT_AMBER

    # Save presentation
    output_path = "exoplanet_hunters_presentation.pptx"
    prs.save(output_path)
    print(f"[OK] PPTX successfully compiled and saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
